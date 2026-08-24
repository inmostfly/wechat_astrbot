"""Tests for the local common-document extraction MCP."""

from __future__ import annotations

from io import BytesIO
import json
import unittest

from docx import Document
from openpyxl import Workbook
from pptx import Presentation

from document_mcp_client import extract_document
from document_mcp_server import DocumentToolError, extract_document_bytes


class DocumentMCPTests(unittest.TestCase):
    def test_utf8_text_and_json_are_extracted(self) -> None:
        text_result = extract_document_bytes(
            "说明.txt",
            "第一行\n第二行".encode("utf-8"),
        )
        json_result = extract_document_bytes(
            "data.json",
            json.dumps({"课程": "人工智能"}, ensure_ascii=False).encode("utf-8"),
        )

        self.assertEqual(text_result["content"], "第一行\n第二行")
        self.assertIn('"课程": "人工智能"', json_result["content"])

    def test_html_ignores_script_content(self) -> None:
        result = extract_document_bytes(
            "page.html",
            b"<h1>Title</h1><script>ignore()</script><p>Body</p>",
        )

        self.assertIn("Title", result["content"])
        self.assertIn("Body", result["content"])
        self.assertNotIn("ignore", result["content"])

    def test_docx_tables_are_extracted(self) -> None:
        document = Document()
        document.add_paragraph("课程安排")
        table = document.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "日期"
        table.cell(0, 1).text = "内容"
        table.cell(1, 0).text = "周一"
        table.cell(1, 1).text = "机器学习"
        output = BytesIO()
        document.save(output)

        result = extract_document_bytes("课程.docx", output.getvalue())

        self.assertIn("课程安排", result["content"])
        self.assertIn("周一 机器学习", result["content"])

    def test_xlsx_sheets_are_extracted(self) -> None:
        workbook = Workbook()
        worksheet = workbook.active
        worksheet.title = "成绩"
        worksheet.append(["姓名", "分数"])
        worksheet.append(["Alice", 95])
        output = BytesIO()
        workbook.save(output)

        result = extract_document_bytes("成绩.xlsx", output.getvalue())

        self.assertIn("[工作表：成绩]", result["content"])
        self.assertIn("Alice 95", result["content"])

    def test_pptx_slides_are_extracted(self) -> None:
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "项目汇报"
        slide.placeholders[1].text = "第一阶段完成"
        output = BytesIO()
        presentation.save(output)

        result = extract_document_bytes("汇报.pptx", output.getvalue())

        self.assertIn("[第1页]", result["content"])
        self.assertIn("项目汇报", result["content"])
        self.assertIn("第一阶段完成", result["content"])

    def test_content_is_truncated_and_legacy_office_is_rejected(self) -> None:
        result = extract_document_bytes("long.txt", b"a" * 2_000, max_chars=1_000)

        self.assertTrue(result["truncated"])
        self.assertIn("后续已截断", result["content"])
        with self.assertRaises(DocumentToolError):
            extract_document_bytes("legacy.doc", b"legacy")

    def test_stdio_mcp_roundtrip(self) -> None:
        result = extract_document(
            "mcp.md",
            "通过 MCP 提取的内容".encode("utf-8"),
            max_chars=2_000,
        )

        self.assertEqual(result["filename"], "mcp.md")
        self.assertEqual(result["content"], "通过 MCP 提取的内容")


if __name__ == "__main__":
    unittest.main()
